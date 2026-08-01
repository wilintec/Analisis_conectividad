#!/usr/bin/env python3
"""Construye un sitio estático de GitHub Pages a partir de un PowerPoint.

El PowerPoint es la fuente maestra. El script:
1. extrae títulos, preguntas, técnicas y texto con python-pptx;
2. renderiza las diapositivas mediante LibreOffice y Poppler;
3. genera imágenes WebP, miniaturas, un índice buscable y HTML responsive;
4. copia el PowerPoint original para descarga.

Diseñado para ejecutarse localmente o en GitHub Actions.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import unicodedata
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable

import yaml
from jinja2 import Environment, FileSystemLoader, select_autoescape
from PIL import Image
from pptx import Presentation


@dataclass
class SlideData:
    number: int
    title: str
    question: str
    technique: str
    summary: str
    text: str
    image: str
    thumbnail: str
    chapter_id: str
    chapter_title: str
    embed: str | None = None


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--config", default="site.yml", help="Archivo YAML de configuración")
    parser.add_argument("--pptx", help="Sobrescribe la ruta del PowerPoint definida en site.yml")
    parser.add_argument("--output", default="dist", help="Directorio de salida")
    parser.add_argument("--keep-build", action="store_true", help="Conservar archivos temporales de render")
    return parser.parse_args()


def run(command: list[str], *, cwd: Path | None = None) -> None:
    printable = " ".join(str(part) for part in command)
    print(f"[build] {printable}")
    completed = subprocess.run(command, cwd=cwd, text=True, capture_output=True)
    if completed.returncode != 0:
        raise RuntimeError(
            f"Comando falló ({completed.returncode}): {printable}\n"
            f"STDOUT:\n{completed.stdout}\nSTDERR:\n{completed.stderr}"
        )


def require_command(name: str) -> str:
    executable = shutil.which(name)
    if not executable:
        raise RuntimeError(
            f"No se encontró '{name}'. En Ubuntu instale LibreOffice y Poppler; "
            "en GitHub Actions el workflow incluido realiza esta instalación."
        )
    return executable


def load_config(path: Path) -> dict[str, Any]:
    if not path.exists():
        raise FileNotFoundError(f"No existe el archivo de configuración: {path}")
    with path.open("r", encoding="utf-8") as fh:
        data = yaml.safe_load(fh) or {}
    return data


def slugify(value: str) -> str:
    normalized = unicodedata.normalize("NFKD", value)
    ascii_text = normalized.encode("ascii", "ignore").decode("ascii")
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", ascii_text).strip("-").lower()
    return slug or "seccion"


def clean_text(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def ordered_texts(slide: Any) -> list[str]:
    """Devuelve textos ordenados de arriba a abajo y de izquierda a derecha."""
    candidates: list[tuple[int, int, str]] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = clean_text(shape.text)
        if not text:
            continue
        candidates.append((int(shape.top), int(shape.left), text))
    candidates.sort(key=lambda item: (item[0], item[1]))
    return [text for _, _, text in candidates]


def extract_title(slide: Any, texts: list[str], number: int) -> str:
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            candidate = clean_text(title_shape.text)
            if candidate:
                if candidate.upper().startswith("CAPÍTULO") and len(texts) > 1:
                    return texts[1]
                return candidate
    except Exception:
        pass
    if texts:
        if texts[0].upper().startswith("CAPÍTULO") and len(texts) > 1:
            return texts[1]
        return texts[0]
    return f"Diapositiva {number}"


def find_after_label(texts: list[str], labels: Iterable[str]) -> str:
    label_set = tuple(label.upper() for label in labels)
    for index, text in enumerate(texts):
        upper = text.upper()
        if any(label in upper for label in label_set):
            # Soporta etiqueta y valor en el mismo cuadro de texto.
            for label in label_set:
                pos = upper.find(label)
                if pos >= 0:
                    tail = clean_text(text[pos + len(label):].lstrip(" :–—-"))
                    if tail:
                        return tail
            if index + 1 < len(texts):
                return texts[index + 1]
    return ""


def infer_summary(texts: list[str], title: str, question: str, technique: str) -> str:
    excluded = {
        title,
        question,
        technique,
        "Conectividad, inclusión y desigualdad digital",
        "PREGUNTA DE INVESTIGACIÓN",
    }
    candidates: list[str] = []
    for text in texts:
        if text in excluded or text.isdigit() or len(text) < 30:
            continue
        if text.startswith("*"):
            continue
        candidates.append(text)
    if candidates:
        return max(candidates, key=len)[:360]
    if question:
        return f"Esta diapositiva responde: {question}"
    return "Contenido extraído del documento maestro."


def chapter_for_slide(number: int, chapter_config: list[dict[str, Any]]) -> tuple[str, str]:
    for item in chapter_config:
        start = int(item.get("start", number))
        end = int(item.get("end", start))
        if start <= number <= end:
            title = clean_text(str(item.get("title", f"Capítulo {number}")))
            return str(item.get("id", slugify(title))), title
    return "contenido", "Contenido"


def extract_slides(
    pptx_path: Path,
    chapter_config: list[dict[str, Any]],
    overrides: dict[str | int, dict[str, Any]],
) -> list[SlideData]:
    presentation = Presentation(str(pptx_path))
    slides: list[SlideData] = []
    for number, slide in enumerate(presentation.slides, start=1):
        texts = ordered_texts(slide)
        title = extract_title(slide, texts, number)
        question = find_after_label(texts, ["PREGUNTA DE INVESTIGACIÓN"])
        technique = find_after_label(texts, ["TÉCNICA ESTADÍSTICA", "TÉCNICA:", "TÉCNICA "])
        summary = infer_summary(texts, title, question, technique)
        chapter_id, chapter_title = chapter_for_slide(number, chapter_config)
        override = overrides.get(number) or overrides.get(str(number)) or {}
        slides.append(
            SlideData(
                number=number,
                title=clean_text(str(override.get("title", title))),
                question=clean_text(str(override.get("question", question))),
                technique=clean_text(str(override.get("technique", technique))),
                summary=clean_text(str(override.get("summary", summary))),
                text="\n".join(texts),
                image="",
                thumbnail="",
                chapter_id=chapter_id,
                chapter_title=chapter_title,
                embed=override.get("embed"),
            )
        )
    return slides


def find_rendered_pdf(render_dir: Path, pptx_path: Path) -> Path:
    preferred = render_dir / f"{pptx_path.stem}.pdf"
    if preferred.exists():
        return preferred
    pdfs = list(render_dir.glob("*.pdf"))
    if len(pdfs) == 1:
        return pdfs[0]
    raise RuntimeError(f"LibreOffice no generó el PDF esperado en {render_dir}")


def render_pptx_to_pngs(pptx_path: Path, build_dir: Path, dpi: int) -> list[Path]:
    soffice = require_command("soffice")
    pdftoppm = require_command("pdftoppm")
    render_dir = build_dir / "render"
    render_dir.mkdir(parents=True, exist_ok=True)
    run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(render_dir), str(pptx_path)])
    pdf_path = find_rendered_pdf(render_dir, pptx_path)
    prefix = render_dir / "slide"
    run([pdftoppm, "-png", "-r", str(dpi), str(pdf_path), str(prefix)])
    pngs = sorted(render_dir.glob("slide-*.png"), key=lambda p: int(p.stem.rsplit("-", 1)[1]))
    if not pngs:
        raise RuntimeError("Poppler no generó imágenes PNG")
    return pngs


def convert_images(
    pngs: list[Path],
    output: Path,
    image_format: str,
    quality: int,
    thumbnail_width: int,
    thumbnail_format: str = "jpg",
) -> tuple[list[str], list[str]]:
    slides_dir = output / "assets" / "slides"
    thumbs_dir = output / "assets" / "thumbs"
    slides_dir.mkdir(parents=True, exist_ok=True)
    thumbs_dir.mkdir(parents=True, exist_ok=True)
    images: list[str] = []
    thumbs: list[str] = []
    extension = image_format.lower().lstrip(".")
    thumb_extension = thumbnail_format.lower().lstrip(".")
    if extension == "jpeg":
        extension = "jpg"
    if thumb_extension == "jpeg":
        thumb_extension = "jpg"

    for number, png_path in enumerate(pngs, start=1):
        image_rel = f"assets/slides/slide-{number:03d}.{extension}"
        thumb_rel = f"assets/thumbs/slide-{number:03d}.{thumb_extension}"
        image_target = output / image_rel
        thumb_target = output / thumb_rel

        # Copiar el PNG renderizado evita una recompresión costosa y mantiene texto nítido.
        if extension == "png":
            shutil.copy2(png_path, image_target)
        else:
            with Image.open(png_path) as image:
                image = image.convert("RGB")
                if extension == "webp":
                    image.save(image_target, "WEBP", quality=quality, method=0)
                elif extension == "jpg":
                    image.save(image_target, "JPEG", quality=quality, subsampling=0)
                else:
                    raise ValueError(f"Formato de imagen no soportado: {extension}")

        with Image.open(png_path) as image:
            image = image.convert("RGB")
            height = max(1, round(image.height * thumbnail_width / image.width))
            image.thumbnail((thumbnail_width, height), Image.Resampling.LANCZOS)
            if thumb_extension == "jpg":
                image.save(thumb_target, "JPEG", quality=min(84, quality), subsampling=1)
            elif thumb_extension == "webp":
                image.save(thumb_target, "WEBP", quality=min(82, quality), method=0)
            elif thumb_extension == "png":
                image.save(thumb_target, "PNG")
            else:
                raise ValueError(f"Formato de miniatura no soportado: {thumb_extension}")

        images.append(image_rel)
        thumbs.append(thumb_rel)
    return images, thumbs


def copy_static_files(repo_root: Path, output: Path) -> None:
    source = repo_root / "static"
    if not source.exists():
        raise FileNotFoundError(f"No existe el directorio static: {source}")
    target = output / "assets"
    shutil.copytree(source, target, dirs_exist_ok=True)


def build_chapters(slides: list[SlideData], config: list[dict[str, Any]]) -> list[dict[str, Any]]:
    chapters: list[dict[str, Any]] = []
    seen: set[str] = set()
    for item in config:
        chapter_id = str(item.get("id", slugify(str(item.get("title", "capitulo")))))
        chapter_slides = [slide for slide in slides if slide.chapter_id == chapter_id]
        if chapter_slides:
            chapters.append({
                "id": chapter_id,
                "title": str(item.get("title", chapter_slides[0].chapter_title)),
                "slides": chapter_slides,
            })
            seen.add(chapter_id)
    for slide in slides:
        if slide.chapter_id not in seen:
            chapters.append({"id": slide.chapter_id, "title": slide.chapter_title, "slides": []})
            seen.add(slide.chapter_id)
        next(ch for ch in chapters if ch["id"] == slide.chapter_id)["slides"].append(slide)
    # Evita duplicación cuando las diapositivas ya fueron añadidas en el primer bucle.
    for chapter in chapters:
        unique: dict[int, SlideData] = {s.number: s for s in chapter["slides"]}
        chapter["slides"] = [unique[key] for key in sorted(unique)]
    return chapters


def write_json(path: Path, payload: Any) -> None:
    with path.open("w", encoding="utf-8") as fh:
        json.dump(payload, fh, ensure_ascii=False, indent=2)


def main() -> int:
    args = parse_args()
    repo_root = Path.cwd().resolve()
    config_path = (repo_root / args.config).resolve()
    config = load_config(config_path)
    site = config.get("site", {})
    presentation_config = config.get("presentation", {})
    chapter_config = config.get("chapters", []) or []
    overrides = config.get("slide_overrides", {}) or {}

    configured_pptx = args.pptx or presentation_config.get("file")
    if not configured_pptx:
        raise ValueError("Defina presentation.file en site.yml o use --pptx")
    pptx_path = Path(configured_pptx)
    if not pptx_path.is_absolute():
        pptx_path = (repo_root / pptx_path).resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(f"No existe el PowerPoint: {pptx_path}")

    output = Path(args.output)
    if not output.is_absolute():
        output = (repo_root / output).resolve()
    if output.exists():
        shutil.rmtree(output)
    output.mkdir(parents=True)

    slides = extract_slides(pptx_path, chapter_config, overrides)

    build_parent = repo_root / ".build"
    if build_parent.exists() and not args.keep_build:
        shutil.rmtree(build_parent)
    build_parent.mkdir(parents=True, exist_ok=True)
    pngs = render_pptx_to_pngs(
        pptx_path,
        build_parent,
        int(presentation_config.get("render_dpi", 144)),
    )
    if len(pngs) != len(slides):
        raise RuntimeError(
            f"El PPTX contiene {len(slides)} diapositivas, pero se renderizaron {len(pngs)} páginas."
        )

    image_paths, thumb_paths = convert_images(
        pngs,
        output,
        str(presentation_config.get("image_format", "webp")),
        int(presentation_config.get("image_quality", 88)),
        int(presentation_config.get("thumbnail_width", 520)),
        str(presentation_config.get("thumbnail_format", "jpg")),
    )
    for slide, image_path, thumb_path in zip(slides, image_paths, thumb_paths, strict=True):
        slide.image = image_path
        slide.thumbnail = thumb_path

    copy_static_files(repo_root, output)
    downloads = output / "downloads"
    downloads.mkdir(exist_ok=True)
    download_name = str(presentation_config.get("download_name", pptx_path.name))
    shutil.copy2(pptx_path, downloads / download_name)

    # Copia embeds opcionales si existen dentro de static; la ruta ya queda disponible bajo assets/.
    chapters = build_chapters(slides, chapter_config)
    generated_at = datetime.now(timezone.utc).astimezone().isoformat(timespec="seconds")

    env = Environment(
        loader=FileSystemLoader(str(repo_root / "templates")),
        autoescape=select_autoescape(["html", "xml"]),
        trim_blocks=True,
        lstrip_blocks=True,
    )
    template = env.get_template("index.html.j2")
    html_output = template.render(
        site=site,
        slides=slides,
        chapters=chapters,
        download_name=download_name,
        generated_at=generated_at,
        slide_dicts=[asdict(slide) for slide in slides],
    )
    (output / "index.html").write_text(html_output, encoding="utf-8")
    (output / "404.html").write_text(html_output, encoding="utf-8")
    (output / ".nojekyll").write_text("", encoding="utf-8")

    manifest = {
        "site": site,
        "generated_at": generated_at,
        "source": pptx_path.name,
        "slide_count": len(slides),
        "slides": [asdict(slide) for slide in slides],
    }
    write_json(output / "manifest.json", manifest)
    write_json(
        output / "search-index.json",
        [
            {
                "number": slide.number,
                "title": slide.title,
                "question": slide.question,
                "technique": slide.technique,
                "summary": slide.summary,
                "text": slide.text,
                "chapter_id": slide.chapter_id,
            }
            for slide in slides
        ],
    )

    if not args.keep_build:
        shutil.rmtree(build_parent, ignore_errors=True)

    print(f"[build] Sitio generado en: {output}")
    print(f"[build] Diapositivas procesadas: {len(slides)}")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:
        print(f"[error] {exc}", file=sys.stderr)
        raise
