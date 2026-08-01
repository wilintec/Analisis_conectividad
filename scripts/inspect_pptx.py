#!/usr/bin/env python3
"""Genera un inventario CSV/JSON del contenido textual del PowerPoint maestro."""

from __future__ import annotations

import argparse
import csv
import json
import re
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation


def clean(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def ordered_texts(slide: Any) -> list[str]:
    values: list[tuple[int, int, str]] = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = clean(shape.text)
            if text:
                values.append((int(shape.top), int(shape.left), text))
    values.sort(key=lambda item: (item[0], item[1]))
    return [item[2] for item in values]


def after_label(texts: list[str], label: str) -> str:
    upper_label = label.upper()
    for index, text in enumerate(texts):
        upper = text.upper()
        if upper_label in upper:
            tail = clean(text[upper.find(upper_label) + len(label):].lstrip(" :–—-"))
            if tail:
                return tail
            if index + 1 < len(texts):
                return texts[index + 1]
    return ""


def title_for(slide: Any, texts: list[str], number: int) -> str:
    try:
        title = clean(slide.shapes.title.text) if slide.shapes.title else ""
    except Exception:
        title = ""
    if title.upper().startswith("CAPÍTULO") and len(texts) > 1:
        return texts[1]
    return title or (texts[0] if texts else f"Diapositiva {number}")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--config", default="site.yml")
    parser.add_argument("--pptx")
    parser.add_argument("--output", default="reports/slide_inventory.csv")
    parser.add_argument("--json-output", default="reports/slide_inventory.json")
    args = parser.parse_args()

    root = Path.cwd()
    config = yaml.safe_load((root / args.config).read_text(encoding="utf-8"))
    configured = args.pptx or config["presentation"]["file"]
    pptx_path = Path(configured)
    if not pptx_path.is_absolute():
        pptx_path = root / pptx_path

    prs = Presentation(str(pptx_path))
    rows: list[dict[str, Any]] = []
    for number, slide in enumerate(prs.slides, 1):
        texts = ordered_texts(slide)
        rows.append({
            "slide": number,
            "title": title_for(slide, texts, number),
            "technique": after_label(texts, "TÉCNICA ESTADÍSTICA") or after_label(texts, "TÉCNICA"),
            "research_question": after_label(texts, "PREGUNTA DE INVESTIGACIÓN"),
            "text_blocks": len(texts),
            "characters": sum(len(text) for text in texts),
            "all_text": " | ".join(texts),
        })

    output = root / args.output
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("w", newline="", encoding="utf-8-sig") as fh:
        writer = csv.DictWriter(fh, fieldnames=rows[0].keys())
        writer.writeheader()
        writer.writerows(rows)

    json_output = root / args.json_output
    json_output.parent.mkdir(parents=True, exist_ok=True)
    json_output.write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Inventario generado: {output}")
    print(f"Inventario JSON: {json_output}")
    print(f"Diapositivas: {len(rows)}")


if __name__ == "__main__":
    main()
